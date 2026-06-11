"""
SafeWave-AI v0.0.3 — Redis DB 데이터 설계 발표 PPT 생성 스크립트
출력: C:/rp5/docs/safewave_redis_db_presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── 색상 팔레트 (db_spec_v2.html 일치) ──────────────────────────────────
BG       = RGBColor(0x0F, 0x16, 0x29)
SURFACE  = RGBColor(0x1E, 0x2D, 0x4A)
SURFACE2 = RGBColor(0x24, 0x33, 0x52)
BORDER   = RGBColor(0x2E, 0x40, 0x70)
BLUE     = RGBColor(0x3B, 0x82, 0xF6)
BLUE_L   = RGBColor(0x93, 0xC5, 0xFD)
TEAL     = RGBColor(0x0E, 0xA5, 0xE9)
GREEN    = RGBColor(0x10, 0xB9, 0x81)
ORANGE   = RGBColor(0xF9, 0x73, 0x16)
PURPLE   = RGBColor(0xA8, 0x55, 0xF7)
PINK     = RGBColor(0xEC, 0x48, 0x99)
YELLOW   = RGBColor(0xEA, 0xB3, 0x08)
WHITE    = RGBColor(0xE2, 0xE8, 0xF0)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
DIM      = RGBColor(0x64, 0x74, 0x8B)
BLACK    = RGBColor(0x05, 0x08, 0x14)


# ── 헬퍼 함수 ────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, x, y, w, h,
                size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
                italic=False, font="Calibri"):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txb


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=0):
    from pptx.util import Pt as _Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE = 1
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_rect_text(slide, text, x, y, w, h, fill_color, text_color=WHITE,
                  size=13, bold=False, line_color=None, line_width=0,
                  align=PP_ALIGN.CENTER, valign="middle", font="Calibri"):
    shape = add_rect(slide, x, y, w, h, fill_color, line_color, line_width)
    tf = shape.text_frame
    tf.word_wrap = True
    # vertical align
    from pptx.enum.text import MSO_ANCHOR
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "top":
        tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    run.font.name = font
    return shape


def add_notes(slide, text):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text


def add_table(slide, data, x, y, w, h,
              header_fill=SURFACE2, header_color=WHITE,
              row_fill=SURFACE, alt_fill=None,
              cell_color=WHITE, border_color=BORDER,
              col_widths=None, header_size=11, cell_size=10):
    """data: list of rows (first row = header if header=True)"""
    rows = len(data)
    cols = len(data[0])
    tbl = slide.shapes.add_table(rows, cols,
                                  Inches(x), Inches(y),
                                  Inches(w), Inches(h))
    table = tbl.table

    # Set column widths
    if col_widths:
        total_w = Inches(w)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = int(total_w * cw)

    for r_idx, row in enumerate(data):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = cell_text

            # Fill
            fill = cell.fill
            fill.solid()
            if r_idx == 0:
                fill.fore_color.rgb = header_fill
            elif alt_fill and r_idx % 2 == 0:
                fill.fore_color.rgb = alt_fill
            else:
                fill.fore_color.rgb = row_fill

            # Text
            tf = cell.text_frame
            tf.word_wrap = True
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.LEFT
                for run in para.runs:
                    run.font.size = Pt(header_size if r_idx == 0 else cell_size)
                    run.font.bold = (r_idx == 0)
                    run.font.color.rgb = header_color if r_idx == 0 else cell_color
                    run.font.name = "Calibri"

    return tbl


def add_arrow(slide, x1, y1, x2, y2, color=MUTED, width=2):
    """Simple horizontal arrow using line connector"""
    from pptx.util import Pt as _Pt
    # Draw as a line shape
    cx = min(x1, x2)
    cy = min(y1, y2) - 0.02
    cw = abs(x2 - x1)
    ch = abs(y2 - y1) + 0.04
    if cw < 0.01:
        cw = 0.01
    conn = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    return conn


# ── 프레젠테이션 초기화 ────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = 13.333
H = 7.5

BLANK_LAYOUT = prs.slide_layouts[6]  # blank


# ════════════════════════════════════════════════════════════════════════
# SLIDE 1 — 타이틀
# ════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s1, BG)

# 상단 컬러 밴드 (subtle teal accent)
add_rect(s1, 0, 0, W, 0.08, TEAL)

# 메인 타이틀
add_textbox(s1, "SafeWave-AI", 1.0, 1.6, 11, 1.6,
            size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

# 서브타이틀
add_textbox(s1, "Redis DB 데이터 설계 — 독거인 안전 모니터링 시스템",
            1.0, 3.4, 11, 0.7, size=22, color=MUTED, align=PP_ALIGN.CENTER)

# 버전 배지
add_rect_text(s1, "v0.0.3", 5.6, 4.3, 1.0, 0.4,
              fill_color=SURFACE2, text_color=BLUE_L, size=14, bold=True,
              line_color=BLUE, line_width=1)

# 팀명
add_textbox(s1, "앰비언트컴퓨팅개발 1  ·  04팀",
            1.0, 5.2, 11, 0.5, size=14, color=DIM, align=PP_ALIGN.CENTER)

# 하단 밴드
add_rect(s1, 0, H - 0.08, W, 0.08, BLUE)

add_notes(s1, """[발표 대본 - Slide 1: 타이틀]

안녕하세요. 저희는 앰비언트컴퓨팅개발 1 04팀입니다.
오늘 발표는 SafeWave-AI 버전 0.0.3의 데이터 설계, 특히 Redis를 중심으로 한 전체 데이터 흐름에 대해 설명 드리겠습니다.

SafeWave-AI는 Raspberry Pi 5 위에서 동작하는 독거인 안전 모니터링 시스템입니다.
WiFi CSI 신호와 마이크 음향을 실시간으로 분석해 낙상, 생체신호 이상, 환경음 위험 상황을 감지하고 보호자에게 알림을 전달합니다.
이 모든 데이터의 중심에 Redis가 있습니다.""")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 2 — 전체 데이터 파이프라인
# ════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s2, BG)

add_textbox(s2, "전체 데이터 파이프라인", 0.5, 0.25, 12, 0.7,
            size=28, bold=True, color=WHITE)
add_textbox(s2, "ESP32-S3 → Redis → AI → 사용자 알림까지의 전체 흐름",
            0.5, 0.9, 12, 0.4, size=14, color=MUTED)

# ── CSI 경로 (상단) ─────────────────────────────────────────────────────
boxes_y = 2.0
BOX_H = 1.1

# ESP32-S3
add_rect_text(s2, "ESP32-S3\nCSI 센서 (1~6대)",
              0.4, boxes_y, 2.1, BOX_H, ORANGE, WHITE, 12, True)
# Arrow
add_arrow(s2, 2.5, boxes_y + BOX_H/2, 3.1, boxes_y + BOX_H/2, ORANGE, 2)
add_textbox(s2, "UDP:5005\n788B", 2.45, boxes_y + 0.55, 0.75, 0.5, size=9, color=ORANGE)

# sensing
add_rect_text(s2, "sensing\n컨테이너",
              3.1, boxes_y, 1.7, BOX_H, SURFACE2, WHITE, 12, True,
              line_color=BORDER, line_width=1)
# Arrow
add_arrow(s2, 4.8, boxes_y + BOX_H/2, 5.4, boxes_y + BOX_H/2, BLUE_L, 2)

# Redis csi:raw
add_rect_text(s2, "Redis\ncsi:raw",
              5.4, boxes_y, 1.7, BOX_H, SURFACE, WHITE, 12, True,
              line_color=BLUE, line_width=1.5)
# Arrow
add_arrow(s2, 7.1, boxes_y + BOX_H/2, 7.7, boxes_y + BOX_H/2, BLUE_L, 2)

# ai engine
add_rect_text(s2, "AI 엔진\nM1~M5 추론",
              7.7, boxes_y, 1.9, BOX_H, SURFACE2, WHITE, 12, True,
              line_color=BORDER, line_width=1)
# Arrow
add_arrow(s2, 9.6, boxes_y + BOX_H/2, 10.2, boxes_y + BOX_H/2, GREEN, 2)

# ai:result
add_rect_text(s2, "Redis\nai:result",
              10.2, boxes_y, 1.7, BOX_H, SURFACE, WHITE, 12, True,
              line_color=GREEN, line_width=1.5)

# ── 오디오 경로 (하단) ──────────────────────────────────────────────────
audio_y = 3.8
add_rect_text(s2, "마이크\nVAD 수집",
              0.4, audio_y, 2.1, BOX_H * 0.85, PURPLE, WHITE, 12, True)
add_arrow(s2, 2.5, audio_y + (BOX_H*0.85)/2, 5.4, audio_y + (BOX_H*0.85)/2, PURPLE, 2)
add_textbox(s2, "audio-sensing", 3.0, audio_y + 0.4, 2.0, 0.4, size=9, color=PURPLE)
add_rect_text(s2, "Redis\naudio:events",
              5.4, audio_y, 1.7, BOX_H * 0.85, SURFACE, WHITE, 12, True,
              line_color=PURPLE, line_width=1.5)
# 오디오 → ai engine (수직 연결)
add_arrow(s2, 6.27, audio_y, 8.65, boxes_y + BOX_H, PURPLE, 1)

# ── 하단 출력 ────────────────────────────────────────────────────────
out_y = 5.3
add_rect_text(s2, "api:8000\nREST / WebSocket", 7.7, out_y, 1.9, 0.9, SURFACE2, BLUE_L, 11, True, line_color=BORDER, line_width=1)
add_rect_text(s2, "MQTT\nsafewave/ai/*",        9.7, out_y, 1.8, 0.9, SURFACE2, ORANGE, 11, True, line_color=BORDER, line_width=1)
add_rect_text(s2, "FCM\n푸시 알림",              11.6, out_y, 1.5, 0.9, SURFACE2, PINK,   11, True, line_color=BORDER, line_width=1)

# ai:result → 출력들
add_arrow(s2, 11.05, boxes_y + BOX_H, 8.65, out_y, GREEN, 1)
add_arrow(s2, 11.05, boxes_y + BOX_H, 10.6, out_y, ORANGE, 1)
add_arrow(s2, 11.05, boxes_y + BOX_H, 12.35, out_y, PINK, 1)

add_notes(s2, """[발표 대본 - Slide 2: 전체 데이터 파이프라인]

데이터는 크게 두 경로로 들어옵니다.

첫 번째는 WiFi CSI 경로입니다.
ESP32-S3 센서 노드 최대 6대가 100Hz로 788바이트 고정 패킷을 UDP로 전송하고, sensing 컨테이너가 수신해서 Redis의 csi:raw 스트림에 씁니다.

두 번째는 오디오 경로입니다.
마이크에서 VAD(Voice Activity Detection)로 트리거된 오디오 이벤트가 audio:events 스트림에 저장됩니다.

AI 엔진은 두 스트림을 동시에 읽어 M1부터 M5까지 병렬 추론하고, 결과를 ai:result 스트림에 씁니다.
최종적으로 API 서버, MQTT 브로커, FCM 푸시 알림으로 보호자에게 전달됩니다.

모든 데이터는 Redis 메모리에만 존재하며 디스크에 기록되지 않습니다.""")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Redis 설계 원칙
# ════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s3, BG)

add_textbox(s3, "Redis 설계 원칙", 0.5, 0.25, 12, 0.7,
            size=28, bold=True, color=WHITE)
add_textbox(s3, "3가지 핵심 제약 — 시스템 전체에 일관되게 적용",
            0.5, 0.9, 12, 0.4, size=14, color=MUTED)

cards = [
    ("💾", "메모리 전용",
     "운영 데이터는 Redis에만 존재\nSSD 파일 캐시 · 로컬 DB 절대 금지\n컨테이너 재시작 시 이력 초기화\n(의도적 설계)",
     BLUE),
    ("⏱", "TTL ≤ 3600초",
     "모든 신규 Redis 키에 EXPIRE 필수\n1시간 이내 자동 만료\n메모리 무한 증가 방지\nStream: MAXLEN + approximate trim",
     GREEN),
    ("🔄", "스트림 3축",
     "csi:raw  → 원시 CSI 데이터\nai:result → 통합 추론 스냅샷\nai:emergency → 위험 이벤트\n단방향 파이프라인",
     ORANGE),
]

for i, (icon, title, body, color) in enumerate(cards):
    cx = 0.7 + i * 4.2
    # Card background
    add_rect(s3, cx, 1.55, 3.9, 4.6, SURFACE, line_color=color, line_width=1.5)
    # Top accent bar
    add_rect(s3, cx, 1.55, 3.9, 0.1, color)
    # Icon
    add_textbox(s3, icon, cx + 0.15, 1.7, 0.8, 0.8, size=32)
    # Title
    add_textbox(s3, title, cx + 0.15, 2.5, 3.5, 0.55,
                size=18, bold=True, color=color)
    # Body
    add_textbox(s3, body, cx + 0.15, 3.1, 3.55, 2.8,
                size=13, color=MUTED)

add_notes(s3, """[발표 대본 - Slide 3: Redis 설계 원칙]

시스템 전체를 관통하는 세 가지 Redis 설계 원칙이 있습니다.

첫째, 메모리 전용입니다.
모든 센서 데이터와 AI 추론 결과는 Redis 메모리에만 저장됩니다.
SSD 파일 쓰기나 로컬 데이터베이스는 사용하지 않습니다.
이는 Raspberry Pi 5의 SSD 수명을 보호하고 I/O 병목을 없애기 위한 의도적인 설계입니다.

둘째, TTL 1시간 원칙입니다.
새로 만드는 모든 Redis 키에는 반드시 3600초 이하의 만료 시간을 붙입니다.
이를 통해 메모리가 무한히 늘어나는 것을 방지합니다.

셋째, 스트림 3축 구조입니다.
csi:raw로 원시 데이터가 들어오고, ai:result로 추론 결과가 나가며, ai:emergency로 위험 이벤트가 분기됩니다.
이 단방향 파이프라인이 시스템의 중심입니다.""")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 4 — csi:raw 스트림
# ════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s4, BG)

add_textbox(s4, "csi:raw 스트림", 0.5, 0.2, 6, 0.65,
            size=28, bold=True, color=ORANGE)
add_textbox(s4, "Redis Stream  |  MAXLEN 36,000  |  100Hz",
            0.5, 0.82, 8, 0.4, size=13, color=MUTED)

# ── 좌측: 788B 패킷 구조 ────────────────────────────────────────────────
add_textbox(s4, "ESP32 UDP 패킷 구조 (788 bytes)", 0.4, 1.35, 6, 0.4,
            size=13, bold=True, color=ORANGE)

pkt_data = [
    ["오프셋", "필드", "타입", "Bytes", "설명"],
    ["0",  "magic",      "char[4]",    "4",   '"CSI!" — 유효성 검사'],
    ["4",  "node_id",    "uint8",      "1",   "노드 번호 (1~6)"],
    ["5",  "reserved",   "uint8",      "1",   "= 0"],
    ["6",  "n_samples",  "uint16",     "2",   "= 64"],
    ["8",  "seq_num",    "uint32",     "4",   "단조 증가 시퀀스"],
    ["12", "ts_ms",      "uint32",     "4",   "Unix 하위 32bit ms"],
    ["16", "rssi",       "int16",      "2",   "수신 신호 강도 (dBm)"],
    ["18", "reserved2",  "uint16",     "2",   "= 0"],
    ["20", "block_raw",  "float32×64", "256", "M1: 광대역 peak-norm"],
    ["276","block_resp", "float32×64", "256", "M2: 0.1–0.6 Hz IIR"],
    ["532","block_heart","float32×64", "256", "M2: 0.8–3.0 Hz IIR"],
]
add_table(s4, pkt_data, 0.4, 1.8, 7.2, 4.8,
          header_fill=RGBColor(0x7C, 0x3A, 0x0A),
          col_widths=[0.08, 0.14, 0.16, 0.08, 0.54])

# ── 우측: Redis 필드 ────────────────────────────────────────────────────
add_textbox(s4, "Redis Stream 필드 (csi:raw)", 7.9, 1.35, 5, 0.4,
            size=13, bold=True, color=ORANGE)

redis_data = [
    ["필드명", "값"],
    ["node",       "node_id (1~6)"],
    ["ts_ms",      "Unix ms (uint32)"],
    ["data_raw",   "raw64 bytes — M1 입력"],
    ["data_resp",  "resp64 bytes — M2 호흡"],
    ["data_heart", "heart64 bytes — M2 심박"],
]
add_table(s4, redis_data, 7.9, 1.8, 5.0, 2.7,
          header_fill=RGBColor(0x7C, 0x3A, 0x0A),
          col_widths=[0.42, 0.58])

# Python struct 표시
add_rect(s4, 7.9, 4.65, 5.0, 1.6, BLACK,
         line_color=ORANGE, line_width=1)
add_textbox(s4,
            'struct.Struct(\n  "<4sBBHIIhH192f"\n)\n→ 788 bytes (고정)',
            8.0, 4.7, 4.8, 1.5,
            size=12, color=ORANGE, font="Consolas")

add_notes(s4, """[발표 대본 - Slide 4: csi:raw 스트림]

csi:raw 스트림은 ESP32-S3 노드에서 전송되는 788바이트 고정 패킷을 저장합니다.

패킷 구조를 보면, 20바이트 헤더와 192개의 float32, 총 768바이트의 데이터로 구성됩니다.
헤더에는 "CSI!"라는 매직 값으로 유효성을 검사하고, 노드 번호, 시퀀스 번호, 타임스탬프, RSSI가 포함됩니다.

데이터 부분은 세 블록으로 나뉩니다.
block_raw는 광대역 peak-norm amplitude로 M1 낙상 감지 모델의 입력입니다.
block_resp는 ESP32에서 0.1~0.6Hz Butterworth IIR 필터를 적용한 호흡 대역 신호입니다.
block_heart는 0.8~3.0Hz 필터를 적용한 심박 대역 신호입니다.

중요한 점은 ESP32가 필터링을 직접 수행한다는 것입니다. Raspberry Pi에서는 scipy 없이 ONNX 추론만 수행합니다.

Python에서는 struct.Struct("<4sBBHIIhH192f")로 한 번에 파싱됩니다.""")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ai:result 스트림
# ════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s5, BG)

add_textbox(s5, "ai:result 스트림", 0.5, 0.2, 7, 0.65,
            size=28, bold=True, color=GREEN)
add_textbox(s5, "Redis Stream  |  MAXLEN 36,000  |  AI 추론 통합 스냅샷 (매 CSI 이벤트마다 1건)",
            0.5, 0.82, 12, 0.4, size=13, color=MUTED)

# ── 좌측: 최상위 필드 ────────────────────────────────────────────────────
add_textbox(s5, "최상위 필드", 0.4, 1.4, 6.2, 0.4,
            size=13, bold=True, color=GREEN)
top_data = [
    ["필드", "타입", "설명"],
    ["ts_ms",         "int",    "CSI 트리거 시각 (Unix ms)"],
    ["node_id",       "int",    "소스 노드 번호"],
    ["risk_level",    "string", '"normal" | "warning" | "critical"'],
    ["risk_score",    "float",  "통합 위험도 0.0–1.0"],
    ["emergency",     "bool",   "위험 임계값 초과 여부"],
    ["emergency_score","float", "도메인 가중 위험 점수"],
    ["slm_summary",   "string", "Qwen-0.5B 판단 근거 텍스트"],
]
add_table(s5, top_data, 0.4, 1.85, 6.2, 3.5,
          header_fill=RGBColor(0x05, 0x56, 0x3A),
          col_widths=[0.30, 0.18, 0.52])

# ── 우측: experts 서브 필드 ───────────────────────────────────────────
add_textbox(s5, "experts 서브 객체 (M1~M4)", 7.0, 1.4, 5.9, 0.4,
            size=13, bold=True, color=GREEN)
exp_data = [
    ["experts.*",       "주요 출력 필드"],
    ["fall (M1)",       "fall_risk: 0~1, fall_label"],
    ["vital (M2)",      "heart_rate, breathing_rate, infer_confidence"],
    ["env_sound (M3)",  "env_sound_label, confidence, audio_ts_ms"],
    ["speech_ko (M4)",  "transcript, language, duration_ms"],
]
add_table(s5, exp_data, 7.0, 1.85, 5.9, 2.4,
          header_fill=RGBColor(0x05, 0x56, 0x3A),
          col_widths=[0.38, 0.62])

# emergency_score 도메인 가중치 카드
add_rect(s5, 7.0, 4.45, 5.9, 2.6, SURFACE, line_color=GREEN, line_width=1)
add_textbox(s5, "응급지수 도메인 가중치", 7.15, 4.55, 5.5, 0.4,
            size=12, bold=True, color=GREEN)
weights = [
    ("M1 낙상",  "40%", ORANGE),
    ("M2 생체",  "30%", PINK),
    ("M3 환경음","15%", PURPLE),
    ("M4 음성",  "15%", TEAL),
]
for i, (name, pct, col) in enumerate(weights):
    bx = 7.15 + i * 1.4
    add_rect(s5, bx, 5.05, 1.25, 0.7, SURFACE2, line_color=col, line_width=1)
    add_textbox(s5, name, bx + 0.05, 5.08, 1.15, 0.3, size=10, color=col)
    add_textbox(s5, pct,  bx + 0.05, 5.35, 1.15, 0.35, size=16, bold=True, color=col)
add_textbox(s5, "복합위험(낙상+음성) 발생 시 × 1.2 보정",
            7.15, 5.88, 5.5, 0.4, size=11, color=MUTED)

add_notes(s5, """[발표 대본 - Slide 5: ai:result 스트림]

ai:result 스트림은 AI 엔진이 매 CSI 이벤트마다 생성하는 통합 추론 스냅샷입니다.

최상위에는 위험 수준, 위험 점수, 응급 여부가 있고, M5 Qwen 모델의 자연어 판단 근거도 포함됩니다.

experts 서브 객체 안에 M1부터 M4까지 각 전문가 모델의 결과가 담깁니다.
M1은 낙상 위험 점수, M2는 심박수와 호흡수, M3는 환경음 분류 레이블, M4는 한국어 STT 결과입니다.

하단의 응급지수는 네 도메인의 가중 합산입니다.
낙상이 40%로 가장 중요하고, 생체신호 30%, 환경음과 음성이 각 15%입니다.
낙상과 음성이 동시에 감지되면 복합위험으로 판단해 1.2배 보정을 적용합니다.""")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ai:emergency 스트림
# ════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s6, BG)

add_textbox(s6, "ai:emergency 스트림", 0.5, 0.2, 8, 0.65,
            size=28, bold=True, color=PINK)
add_textbox(s6, "Redis Stream  |  MAXLEN 3,600  |  warning·critical 이벤트만 기록",
            0.5, 0.82, 12, 0.4, size=13, color=MUTED)

# ── 발생 조건 ────────────────────────────────────────────────────────────
add_textbox(s6, "발생 조건", 0.4, 1.4, 6, 0.4, size=13, bold=True, color=PINK)
cond_data = [
    ["조건", "값"],
    ["risk_level",    '"warning" 또는 "critical"'],
    ["기본 임계값",   "risk_score ≥ 0.6"],
    ["임계값 변경",   "POST /settings → sys:settings"],
    ["SLM 쿨다운",   "최소 3초 간격 (SLM_MIN_INTERVAL_MS)"],
]
add_table(s6, cond_data, 0.4, 1.85, 6.0, 2.4,
          header_fill=RGBColor(0x6B, 0x16, 0x3A),
          col_widths=[0.35, 0.65])

# ── 필드 ──────────────────────────────────────────────────────────────
add_textbox(s6, "스트림 필드", 0.4, 4.4, 6, 0.4, size=13, bold=True, color=PINK)
field_data = [
    ["필드", "설명"],
    ["data",         "JSON 직렬화된 ai:result 전체 스냅샷"],
    ["risk_level",   '"warning" | "critical"'],
    ["ts_ms",        "이벤트 발생 시각"],
]
add_table(s6, field_data, 0.4, 4.85, 6.0, 1.9,
          header_fill=RGBColor(0x6B, 0x16, 0x3A),
          col_widths=[0.3, 0.7])

# ── 하류 (Downstream) ──────────────────────────────────────────────────
add_textbox(s6, "Downstream — 위험 이벤트 전달 경로", 7.2, 1.4, 5.7, 0.4,
            size=13, bold=True, color=PINK)

downstreams = [
    ("📱", "FCM 푸시 알림",
     "critical: android_priority=high\niOS CriticalSound\n잠금화면 전면 인텐트",
     PINK),
    ("📡", "MQTT 발행",
     "safewave/ai/emergency\nHome Assistant 자동화\n피드백 루프 연동",
     ORANGE),
    ("🌐", "API / WebSocket",
     "GET /history\nWS /ws/monitor 실시간\n보호자 앱 UI 표시",
     BLUE),
    ("💾", "오디오 클립",
     "ai:clip:{ts_ms} Redis 저장\n전후 15초 이벤트 메타데이터\nGET /emergency/clip/{ts_ms}",
     PURPLE),
]
for i, (icon, title, body, col) in enumerate(downstreams):
    cx = 7.2 + (i % 2) * 2.9
    cy = 1.9 + (i // 2) * 2.5
    add_rect(s6, cx, cy, 2.65, 2.2, SURFACE, line_color=col, line_width=1)
    add_rect(s6, cx, cy, 2.65, 0.08, col)
    add_textbox(s6, icon + " " + title, cx + 0.12, cy + 0.15, 2.4, 0.45,
                size=12, bold=True, color=col)
    add_textbox(s6, body, cx + 0.12, cy + 0.6, 2.4, 1.45, size=11, color=MUTED)

add_notes(s6, """[발표 대본 - Slide 6: ai:emergency 스트림]

ai:emergency는 위험 상황이 감지됐을 때만 기록되는 스트림입니다.
MAXLEN 3,600으로 ai:result보다 작게 유지됩니다.

발생 조건은 risk_score가 기본 임계값 0.6 이상이거나 risk_level이 warning 또는 critical일 때입니다.
임계값은 /settings API로 실시간 변경 가능합니다.

위험 이벤트가 발생하면 네 경로로 동시에 전달됩니다.

첫째, FCM으로 보호자 스마트폰에 high priority 푸시를 보냅니다.
critical 등급이면 iOS Critical Sound로 무음 모드도 무시하고 알립니다.

둘째, MQTT safewave/ai/emergency 토픽으로 발행해 Home Assistant 자동화와 연동됩니다.

셋째, WebSocket을 통해 보호자 앱 UI에 실시간으로 표시됩니다.

넷째, 전후 15초 오디오 이벤트 메타데이터가 Redis에 저장되어 나중에 조회할 수 있습니다.""")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 7 — audio:events 스트림 + M2 시간축 버퍼
# ════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s7, BG)

add_textbox(s7, "audio:events 스트림  +  M2 시간축 버퍼", 0.5, 0.2, 12, 0.65,
            size=26, bold=True, color=PURPLE)
add_textbox(s7, "Redis Stream  |  MAXLEN 3,600  |  VAD 트리거 오디오 이벤트",
            0.5, 0.82, 12, 0.4, size=13, color=MUTED)

# audio:events 필드 테이블
add_textbox(s7, "audio:events 필드", 0.4, 1.4, 6, 0.4,
            size=13, bold=True, color=PURPLE)
audio_data = [
    ["필드", "타입", "설명"],
    ["node",        "int",    "소스 노드 번호"],
    ["ts_ms",       "int",    "VAD 트리거 시각"],
    ["data",        "JSON",   "sample_rate, channels, duration_ms, peak_db, waveform[]"],
]
add_table(s7, audio_data, 0.4, 1.85, 6.5, 2.0,
          header_fill=RGBColor(0x4A, 0x16, 0x6B),
          col_widths=[0.22, 0.14, 0.64])

# audio → M3/M4 파이프라인
add_textbox(s7, "오디오 → M3/M4 파이프라인", 0.4, 4.0, 6.5, 0.4,
            size=13, bold=True, color=PURPLE)
pipe_data = [
    ["윈도우", "모델", "ms"],
    ["M3_AUDIO_WINDOW_MS", "M3 환경음 (AST)", "3,000 ms"],
    ["M4_AUDIO_WINDOW_MS", "M4 STT (Whisper)", "5,000 ms"],
]
add_table(s7, pipe_data, 0.4, 4.45, 6.5, 1.8,
          header_fill=RGBColor(0x4A, 0x16, 0x6B),
          col_widths=[0.42, 0.35, 0.23])

# M2 시간축 버퍼 (우측)
add_rect(s7, 7.2, 1.35, 5.7, 5.6, SURFACE, line_color=PINK, line_width=1.5)
add_rect(s7, 7.2, 1.35, 5.7, 0.1, PINK)
add_textbox(s7, "M2 시간축 누적 버퍼 (신규)", 7.35, 1.5, 5.3, 0.5,
            size=14, bold=True, color=PINK)

add_textbox(s7, "문제: csi:raw의 data_resp·data_heart는\n64개 서브캐리어 '공간 스냅샷'.\n단일 프레임으로 호흡 주파수(0.1Hz) 추출 불가.",
            7.35, 2.05, 5.3, 1.2, size=12, color=RGBColor(0xFC, 0xA5, 0xA5))

add_textbox(s7, "해결: per-node deque(maxlen=300 @ 100Hz)\n→ 64채널 평균(axis=1) → (N,) 시간 시리즈\n→ FFT로 호흡(0.1–0.6Hz) / 심박(0.8–3.0Hz) 추출",
            7.35, 3.35, 5.3, 1.5, size=12, color=RGBColor(0x86, 0xEF, 0xAC))

buf_data = [
    ["파라미터", "기본값"],
    ["M2_CSI_WINDOW_FRAMES", "300 (3초 @ 100Hz)"],
    ["호흡 완전 해상도",      "1000프레임 = 10초"],
    ["심박 최소 윈도우",      "100프레임 = 1초"],
    ["튜닝 방법",            "env var 오버라이드"],
]
add_table(s7, buf_data, 7.35, 5.0, 5.3, 1.7,
          header_fill=RGBColor(0x6B, 0x16, 0x3A),
          col_widths=[0.48, 0.52])

add_notes(s7, """[발표 대본 - Slide 7: audio:events + M2 시간축 버퍼]

audio:events 스트림은 마이크에서 VAD로 감지된 오디오 이벤트를 저장합니다.
data 필드 안에 PCM 파형, 샘플레이트, peak_db 등이 JSON으로 들어있습니다.

AI 엔진은 CSI 이벤트가 올 때마다 가장 최근 오디오 이벤트들을 시간 윈도우로 병합해 M3와 M4에 전달합니다.
M3는 3초, M4는 5초 윈도우를 사용합니다.

우측 패널은 이번에 새로 추가된 M2 시간축 버퍼입니다.

기존 문제는 csi:raw의 data_resp와 data_heart가 64개 서브캐리어의 순간 스냅샷이라는 점입니다.
이 64개 값으로 FFT를 돌리면 공간 주파수가 나오지, 시간 주파수인 호흡수가 나오지 않습니다.

해결책으로 ai/main.py에 per-node deque를 추가했습니다.
300프레임, 즉 100Hz에서 3초치를 누적하고 64채널을 평균 내어 (N,) 형태의 시간 시리즈를 만듭니다.
이 시리즈에 FFT를 적용하면 올바른 호흡·심박 주파수가 추출됩니다.""")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Expert Latest Keys (ai:m1~m4:latest)
# ════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s8, BG)

add_textbox(s8, "전문가 최신 키  ai:m1:latest ~ ai:m4:latest", 0.5, 0.2, 12, 0.65,
            size=26, bold=True, color=TEAL)
add_textbox(s8, "Redis String  |  TTL 3600s  |  전문가별 최신 추론 결과 스냅샷",
            0.5, 0.82, 12, 0.4, size=13, color=MUTED)

experts = [
    ("M1", "ai:m1:latest", "낙상 감지\n(WiFi Pose)", ORANGE,
     ["fall_risk: float 0–1", "fall_label: str", "infer_source: 'onnx'|'heuristic'",
      "infer_latency_ms: float", "node_id, ts_ms"]),
    ("M2", "ai:m2:latest", "생체신호\n(Frenel Vital)", PINK,
     ["heart_rate: float (bpm)", "breathing_rate: float (bpm)", "infer_confidence: 0–0.65",
      "infer_source: 'onnx'|'fft'", "node_id, ts_ms"]),
    ("M3", "ai:m3:latest", "환경음 분류\n(AST)", PURPLE,
     ["env_sound_label: 7종", "env_sound_confidence: float", "ast_top_class: str",
      "audio_ts_ms, duration_ms", "infer_source: 'onnx'|'heuristic'"]),
    ("M4", "ai:m4:latest", "한국어 STT\n(Whisper)", BLUE,
     ["transcript: str", "language: 'ko'|'unknown'", "duration_ms: float",
      "waveform_rms: float", "infer_source: 'onnx'|'no-audio'"]),
]

for i, (mid, key, title, col, fields) in enumerate(experts):
    cx = 0.4 + i * 3.2
    # Card
    add_rect(s8, cx, 1.5, 3.0, 5.3, SURFACE, line_color=col, line_width=1.5)
    add_rect(s8, cx, 1.5, 3.0, 0.12, col)
    # Model badge
    add_rect_text(s8, mid, cx + 0.12, 1.7, 0.55, 0.45,
                  fill_color=col, text_color=WHITE, size=15, bold=True)
    # Title
    add_textbox(s8, title, cx + 0.75, 1.68, 2.0, 0.5, size=12, bold=True, color=col)
    # Key name
    add_rect(s8, cx + 0.12, 2.27, 2.75, 0.35, BLACK, line_color=col, line_width=0.5)
    add_textbox(s8, key, cx + 0.18, 2.29, 2.6, 0.32, size=9.5, color=col, font="Consolas")
    # Fields
    for j, f in enumerate(fields):
        add_textbox(s8, "▸ " + f, cx + 0.15, 2.78 + j * 0.54, 2.7, 0.48,
                    size=10.5, color=MUTED)

add_notes(s8, """[발표 대본 - Slide 8: Expert Latest Keys]

ai:m1:latest부터 ai:m4:latest까지 네 개의 키는 각 전문가 모델의 가장 최신 추론 결과를 Redis String으로 저장합니다.
TTL은 3600초입니다.

M1은 낙상 위험 점수를 0에서 1 사이로 출력합니다.
M2는 심박수와 호흡수, 그리고 추론 신뢰도를 반환합니다.
방금 설명한 시간축 버퍼 덕분에 FFT 기반 추론이 올바르게 동작합니다.

M3는 silence, speech, music, impact, noise, alarm, unknown의 7가지 환경음을 분류합니다.
낙상을 나타내는 impact나 비명을 포함한 speech가 높으면 위험도가 올라갑니다.

M4는 Whisper 기반 한국어 STT로 "도와줘"나 "아파" 같은 위험 발화를 감지합니다.

각 모델은 ONNX 모델이 있으면 ONNX 추론, 없으면 규칙 기반 폴백으로 동작합니다.
infer_source 필드로 어느 경로로 추론됐는지 추적할 수 있습니다.""")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Node Health & System Keys
# ════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s9, BG)

add_textbox(s9, "보조 키 — 노드 헬스 & 시스템", 0.5, 0.2, 12, 0.65,
            size=26, bold=True, color=YELLOW)
add_textbox(s9, "운영 상태 모니터링 및 시스템 설정 관리",
            0.5, 0.82, 12, 0.4, size=13, color=MUTED)

# 노드 헬스 테이블
add_textbox(s9, "노드 헬스 키", 0.4, 1.4, 6.3, 0.4,
            size=13, bold=True, color=YELLOW)
health_data = [
    ["키 패턴", "타입", "TTL", "설명"],
    ["node:N:last_seen", "String", "30s",   "노드 마지막 수신 Unix 시각"],
    ["node:N:health",    "Hash",   "3600s", "rx / lost / loss_rate / rssi / last_seq"],
]
add_table(s9, health_data, 0.4, 1.85, 6.3, 2.0,
          header_fill=RGBColor(0x6B, 0x56, 0x00),
          col_widths=[0.30, 0.14, 0.12, 0.44])

# 헬스 Hash 필드 상세
add_textbox(s9, "node:N:health Hash 필드 상세", 0.4, 3.95, 6.3, 0.4,
            size=13, bold=True, color=YELLOW)
hash_data = [
    ["필드", "설명"],
    ["rx",         "총 수신 패킷 수"],
    ["lost",       "추정 유실 패킷 수 (시퀀스 갭)"],
    ["loss_rate",  "유실률 (0.0~1.0)"],
    ["rssi",       "마지막 수신 RSSI (dBm)"],
    ["last_seen",  "마지막 수신 Unix 시각"],
    ["last_seq",   "마지막 시퀀스 번호"],
]
add_table(s9, hash_data, 0.4, 4.4, 6.3, 2.75,
          header_fill=RGBColor(0x6B, 0x56, 0x00),
          col_widths=[0.28, 0.72])

# 시스템 키 테이블 (우측)
add_textbox(s9, "시스템 & 기타 키", 7.1, 1.4, 5.8, 0.4,
            size=13, bold=True, color=YELLOW)
sys_data = [
    ["키", "타입", "TTL", "설명"],
    ["sys:settings",       "Hash",   "3600s", "risk_threshold, active_nodes, ai_enabled"],
    ["agg:minute:*",       "Hash",   "3600s", "분 단위 평균 위험도·HR·RR"],
    ["fcm:token:*",        "String", "3600s", "FCM 기기 등록 토큰"],
    ["mqtt:feedback:last", "String", "3600s", "MQTT feedback 마지막 값"],
    ["ai:clip:{ts_ms}",    "String", "3600s", "응급 전후 15초 오디오 메타"],
]
add_table(s9, sys_data, 7.1, 1.85, 5.8, 3.3,
          header_fill=RGBColor(0x6B, 0x56, 0x00),
          col_widths=[0.30, 0.13, 0.12, 0.45])

# API 연계 카드
add_rect(s9, 7.1, 5.35, 5.8, 1.8, SURFACE, line_color=YELLOW, line_width=1)
add_textbox(s9, "API 연계", 7.25, 5.45, 5.4, 0.4, size=12, bold=True, color=YELLOW)
api_links = [
    "GET /nodes/health  →  node:N:health 조회",
    "GET /system/redis-memory  →  INFO memory",
    "POST /settings  →  sys:settings hset + expire",
    "GET /charts/minute  →  agg:minute:* 시각화",
]
for j, line in enumerate(api_links):
    add_textbox(s9, "• " + line, 7.25, 5.88 + j * 0.3, 5.4, 0.32,
                size=10.5, color=MUTED, font="Consolas")

add_notes(s9, """[발표 대본 - Slide 9: Node Health & System Keys]

보조 키들은 운영 모니터링과 시스템 설정에 쓰입니다.

노드 헬스 키는 두 가지입니다.
node:N:last_seen은 TTL 30초로, 이 키가 살아있으면 해당 노드가 온라인임을 의미합니다.
node:N:health Hash에는 수신 패킷 수, 유실 패킷 수, RSSI 등 패킷 통계가 들어있습니다.
GET /nodes/health API로 모든 노드의 상태를 한 번에 조회할 수 있습니다.

시스템 키에서 sys:settings는 위험 임계값, 활성 노드 목록, AI 활성화 여부 같은 런타임 설정입니다.
이 값은 POST /settings로 변경하면 다음 CSI 이벤트부터 즉시 반영됩니다.

agg:minute 키는 분 단위 평균값을 저장해 모니터 대시보드의 추이 차트에 쓰입니다.

ai:clip 키는 응급 이벤트 발생 시 전후 15초 오디오 이벤트 메타데이터를 저장하며, 보호자 앱에서 후속 조회가 가능합니다.""")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Redis 키 전체 맵 (요약)
# ════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s10, BG)

add_textbox(s10, "Redis 키 전체 맵 — 요약", 0.5, 0.2, 12, 0.65,
            size=28, bold=True, color=WHITE)
add_textbox(s10, "SafeWave-AI v0.0.3  |  모든 데이터는 메모리 전용, 컨테이너 재시작 시 이력 없음",
            0.5, 0.82, 12, 0.4, size=13, color=MUTED)

full_data = [
    ["키",                     "타입",   "MAXLEN / TTL", "담당 서비스",  "색상/도메인"],
    ["csi:raw",                "Stream", "MAXLEN 36,000","sensing",      "CSI (주축)"],
    ["audio:events",           "Stream", "MAXLEN 3,600", "audio-sensing","오디오"],
    ["ai:result",              "Stream", "MAXLEN 36,000","ai",           "추론 결과 (주축)"],
    ["ai:emergency",           "Stream", "MAXLEN 3,600", "ai",           "위험 이벤트 (주축)"],
    ["ai:m1:latest",           "String", "TTL 3600s",    "ai",           "M1 낙상 최신"],
    ["ai:m2:latest",           "String", "TTL 3600s",    "ai",           "M2 생체 최신"],
    ["ai:m3:latest",           "String", "TTL 3600s",    "ai",           "M3 환경음 최신"],
    ["ai:m4:latest",           "String", "TTL 3600s",    "ai",           "M4 STT 최신"],
    ["agg:minute:*",           "Hash",   "TTL 3600s",    "ai",           "분 집계 차트"],
    ["node:N:last_seen",       "String", "TTL 30s",      "sensing",      "노드 온라인 여부"],
    ["node:N:health",          "Hash",   "TTL 3600s",    "sensing",      "패킷 통계"],
    ["sys:settings",           "Hash",   "TTL 3600s",    "api",          "런타임 설정"],
    ["fcm:token:*",            "String", "TTL 3600s",    "api",          "FCM 토큰"],
    ["mqtt:feedback:last",     "String", "TTL 3600s",    "ai",           "MQTT 피드백"],
    ["ai:clip:{ts_ms}",        "String", "TTL 3600s",    "ai",           "응급 오디오 클립"],
]

add_table(s10, full_data, 0.4, 1.4, 12.5, 5.7,
          header_fill=SURFACE2,
          row_fill=SURFACE,
          alt_fill=RGBColor(0x18, 0x24, 0x3C),
          col_widths=[0.22, 0.10, 0.17, 0.14, 0.37],
          header_size=11, cell_size=10)

add_notes(s10, """[발표 대본 - Slide 10: Redis 키 전체 맵 요약]

마지막으로 전체 Redis 키를 한 표로 정리했습니다.

핵심 스트림 3개 — csi:raw, ai:result, ai:emergency — 가 시스템의 중심축입니다.
이 세 스트림이 센서 데이터를 AI 추론으로 변환하고 위험 이벤트를 분기하는 파이프라인입니다.

String 타입의 latest 키 4개는 각 전문가 모델의 최신 결과를 즉시 조회할 때 씁니다.

Hash 타입 키들은 분 집계, 노드 헬스, 시스템 설정에 쓰입니다.

모든 키는 TTL 1시간 이내입니다.
스트림은 MAXLEN으로 길이를 제한하고, String/Hash는 EXPIRE로 만료합니다.

node:last_seen만 예외적으로 30초 TTL입니다. 30초가 지나도 키가 살아있으면 해당 노드가 오프라인으로 판단합니다.

이상으로 SafeWave-AI의 Redis 데이터 설계 발표를 마칩니다. 감사합니다.""")


# ════════════════════════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════════════════════════
output_path = r"C:\rp5\docs\safewave_redis_db_presentation.pptx"
prs.save(output_path)
print(f"✅ PPT 저장 완료: {output_path}")
print(f"   슬라이드 수: {len(prs.slides)}")
